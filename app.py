# -*- coding: utf-8 -*-
"""
Streamlit Web App for Universal File-to-Text Conversion.

This script creates a simple web interface where users can:
1. Upload a file (.docx, .xlsx, .pptx, .pdf, .html) via a drag-and-drop interface.
2. The app processes the file into plain text or Markdown.
3. It displays a preview of the first 1000 characters of the extracted text.
4. It provides a download button to save the full text as a markdown file.
5. It shows a rendered preview of the markdown output.
6. It includes a tab to compare original and converted file sizes.
"""

# STEP 1: Import necessary libraries
import streamlit as st
import os
import tempfile
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader
from markdownify import markdownify as md

def convert_file_to_text(filepath):
    """
    Converts the content of a given file to plain text or Markdown.
    This function is adapted from the original Colab script.

    Args:
        filepath (str): The path to the file to be converted.

    Returns:
        str: The extracted text content, or an error string if conversion fails.
    """
    try:
        _, extension = os.path.splitext(filepath)
        extension = extension.lower()
        text_content = ""

        if extension == '.docx':
            document = Document(filepath)
            for paragraph in document.paragraphs:
                text_content += paragraph.text + '\n'
        elif extension == '.xlsx':
            workbook = load_workbook(filepath)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"--- Sheet: {sheet_name} ---\n"
                for row in sheet.iter_rows():
                    row_text = [str(cell.value) if cell.value is not None else "" for cell in row]
                    text_content += "\t".join(row_text) + '\n'
        elif extension == '.pptx':
            presentation = Presentation(filepath)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + '\n'
        elif extension in ['.html', '.htm']:
            with open(filepath, 'r', encoding='utf-8') as f:
                html_content = f.read()
            text_content = md(html_content)
        elif extension == '.pdf':
            reader = PdfReader(filepath)
            for page in reader.pages:
                text_content += page.extract_text() + '\n'
        else:
            return f"Error: Unsupported file type '{extension}'. Please upload a supported file."

        return text_content

    except Exception as e:
        return f"Error processing file: {e}. Please ensure the file is not corrupted."

def format_size(size_bytes):
    """Formats size in bytes to a human-readable string (KB, MB)."""
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024**2:
        return f"{size_bytes/1024:.2f} KB"
    else:
        return f"{size_bytes/1024**2:.2f} MB"

# --- Streamlit App UI ---

st.title("Docs to Markdown Converter")
st.markdown("Drag, drop, and download. It's that simple.")

uploaded_file = st.file_uploader(
    "Choose a file",
    type=['pdf', 'pptx', 'docx', 'xlsx', 'html', 'htm']
)

if uploaded_file is not None:
    if uploaded_file.size > 50 * 1024 * 1024:
        st.warning("Warning: You have uploaded a large file. Processing may take some time.")

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_filepath = os.path.join(temp_dir, uploaded_file.name)
        with open(temp_filepath, "wb") as f:
            f.write(uploaded_file.getbuffer())

        with st.spinner("Converting..."):
            full_text = convert_file_to_text(temp_filepath)

    if not full_text.startswith("Error"):
        st.success("File processed successfully!")

        # Create two tabs: one for the converter output and one for size comparison
        tab1, tab2 = st.tabs(["Converter Output", "File Size Comparison"])

        with tab1:
            st.subheader("Preview (First 1000 characters)")
            st.text_area("Output Preview", full_text[:1000], height=250, disabled=True)

            output_filename = f"converted_{os.path.splitext(uploaded_file.name)[0]}.md"

            st.download_button(
               label="Download Markdown File",
               data=full_text,
               file_name=output_filename,
               mime="text/markdown"
            )

            with st.expander("Rendered Preview"):
                st.markdown(full_text, unsafe_allow_html=True)
        
        with tab2:
            st.subheader("File Size Comparison")
            original_size = uploaded_file.size
            converted_size = len(full_text.encode('utf-8'))
            
            # Create a clean table using markdown
            st.markdown(f"""
            | File Type | Size |
            |-----------|------|
            | Original File | {format_size(original_size)} |
            | Converted Text | {format_size(converted_size)} |
            """)

            if original_size > 0:
                reduction = ((original_size - converted_size) / original_size) * 100
                st.info(f"Text version is **{reduction:.0f}% smaller** than the original file.")

    else:
        st.error(full_text)

